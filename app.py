# app.py
# Future-Ready Workforce (Executive Edition)
# Run: streamlit run app.py

import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from scipy.stats import triang, norm
import streamlit as st
from datetime import datetime
from io import BytesIO

# ----------------------------
# Helper: PPT export (optional)
# ----------------------------
def export_to_ppt(summary_df, hist_data, scenario_name):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    prs = Presentation()
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Future Ready Workforce Simulation"
    slide.placeholders[1].text = f"Scenario: {scenario_name}\nGenerated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    # KPI slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tf = slide.shapes.title.text_frame
    slide.shapes.title.text = "Executive Summary (KPIs)"
    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.9), Inches(4.5)).text_frame
    body.word_wrap = True
    for _, row in summary_df.iterrows():
        p = body.add_paragraph()
        p.text = f"{row['KPI']}: {row['Value']}"
        p.level = 0
        p.font.size = Pt(20)
    # Histogram slide (ROI months)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "ROI Timeline Distribution (Months to Breakeven)"
    fig = px.histogram(pd.DataFrame({"ROI_Months": hist_data}), x="ROI_Months", nbins=40, opacity=0.85)
    fig.update_layout(margin=dict(l=0,r=0,t=30,b=0))
    img = BytesIO()
    fig.write_image(img, format="png", width=1280, height=720, scale=2)
    img.seek(0)
    slide.shapes.add_picture(img, Inches(0.5), Inches(1.5), width=Inches(9.0))
    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out

# ----------------------------
# Styling
# ----------------------------
st.set_page_config(page_title="Future Ready Workforce", page_icon="📊", layout="wide")
st.markdown("""
<style>
.kpi-card {padding:14px;border-radius:16px;background:#0f172a;color:white;box-shadow:0 10px 30px rgba(2,6,23,0.35);}
.kpi-label {font-size:0.9rem;color:#94a3b8;margin-bottom:6px;}
.kpi-value {font-size:1.6rem;font-weight:700;}
.kpi-sub {font-size:0.85rem;color:#cbd5e1;}
.small {font-size:0.85rem;color:#64748b;}
</style>
""", unsafe_allow_html=True)

st.title("📊 Future Ready Workforce")
st.caption("Executive-ready interactive model: productivity target probability, hard savings, satisfaction, and ROI timeline distribution.")

# ----------------------------
# Presets
# ----------------------------
PRESETS = {
    "Pilot program": {
        "participants": 500,
        "budget": 75_000,
        "mgmt_support": "Medium",
        "econ": "Stable",
        "hours_per_emp": 1.0
    },
    "Conservative rollout": {
        "participants": 1000,
        "budget": 150_000,
        "mgmt_support": "Low",
        "econ": "Stable",
        "hours_per_emp": 1.5
    },
    "Aggressive launch": {
        "participants": 3000,
        "budget": 300_000,
        "mgmt_support": "High",
        "econ": "Growth",
        "hours_per_emp": 2.5
    }
}

with st.sidebar:
    st.header("🎛️ Controls")
    scenario = st.selectbox("Preset scenarios", list(PRESETS.keys()) + ["Custom configuration"], index=1)
    runs = st.slider("simulation runs", 2_000, 50_000, 10_000, step=1_000)
    st.markdown("---")
    st.subheader("Core Inputs")

    if scenario != "Custom configuration":
        cfg = PRESETS[scenario]
        participants = st.number_input("Program participants", 100, 50_000, cfg["participants"], step=100)
        hours = st.slider("Monthly time per employee (hrs)", 0.5, 3.0, float(cfg["hours_per_emp"]), 0.1)
        budget = st.number_input("Initial budget (USD)", 50_000, 5_000_000, cfg["budget"], step=5_000)
        mgmt_support = st.selectbox("Management support level", ["Low", "Medium", "High"],
                                    index=["Low","Medium","High"].index(cfg["mgmt_support"]))
        econ = st.selectbox("Economic conditions", ["Recession", "Stable", "Growth"],
                            index=["Recession","Stable","Growth"].index(cfg["econ"]))
    else:
        participants = st.number_input("Program participants", 100, 50_000, 1_000, step=100)
        hours = st.slider("Monthly time per employee (hrs)", 0.5, 3.0, 1.5, 0.1)
        budget = st.number_input("Initial budget (USD)", 50_000, 5_000_000, 150_000, step=5_000)
        mgmt_support = st.selectbox("Management support level", ["Low", "Medium", "High"], index=1)
        econ = st.selectbox("Economic conditions", ["Recession", "Stable", "Growth"], index=1)

    st.subheader("Financial & Workforce Assumptions")
    fully_loaded_hourly = st.number_input("Fully-loaded hourly cost (USD)", 30, 300, 85, step=5)
    baseline_redundant_hours = st.slider("Redundant hours per employee per month", 2, 20, 8, step=1)
    decision_volume = st.number_input("Decisions per month affecting timelines", 100, 100_000, 5_000, step=100)
    avg_hours_per_decision = st.slider("Avg hours across stakeholders per decision", 0.5, 10.0, 2.0, 0.1)

    st.subheader("Churn & Salary")
    baseline_turnover = st.slider("Baseline annual turnover (%)", 5, 40, 20, step=1)
    avg_annual_salary = st.number_input("Average annual salary (USD)", 40_000, 500_000, 120_000, step=5_000)
    replacement_cost_factor = st.slider("Replacement cost factor (x salary)", 0.2, 1.5, 0.5, 0.05)

    st.subheader("Targets & Distributions")
    prod_target = st.slider("Productivity improvement target (%)", 1, 25, 10, step=1)
    decision_reduction_min = st.slider("Decision time reduction min (%)", 5, 40, 10, step=1)
    decision_reduction_max = st.slider("Decision time reduction max (%)", 10, 60, 30, step=1)

    st.markdown("---")
    st.caption("Tip: Export a one-pager to PowerPoint from the bottom of the page.")

# ----------------------------
# Encodings for multipliers
# ----------------------------
mgmt_multiplier_map = {"Low": 0.85, "Medium": 1.00, "High": 1.25}
econ_multiplier_map = {"Recession": 0.85, "Stable": 1.00, "Growth": 1.15}

mgmt_mult = mgmt_multiplier_map[mgmt_support]
econ_mult = econ_multiplier_map[econ]

# ----------------------------
# Random variable generators
# ----------------------------
rng = np.random.default_rng(seed=42)  # stable for demos; could expose as control

def triangular_sample(min_v, mode_v, max_v, size):
    # scipy triang c= (mode - min) / (max - min)
    c = (mode_v - min_v) / (max_v - min_v)
    return triang.rvs(c, loc=min_v, scale=max_v - min_v, size=size, random_state=rng)

def bounded_normal(mean, sd, lower=None, upper=None, size=1):
    x = rng.normal(mean, sd, size)
    if lower is not None: x = np.maximum(x, lower)
    if upper is not None: x = np.minimum(x, upper)
    return x

# ----------------------------
# Model assumptions & formulas
# ----------------------------
N = runs

# 1) Participation variance (who actually engages meaningfully)
#    Assume realized engaged participants ~ Triangular around planned
engaged_participants = np.clip(triangular_sample(
    max(100, participants * 0.7), participants, participants * 1.1, N).astype(int), 1, None)

# 2) Employee time variance (some do more/less than planned)
hours_per_emp = rng.uniform(low=max(0.5, hours*0.7), high=min(3.0, hours*1.2), size=N)

# 3) Budget variance (procurement, vendor variability)
budget_spend = triangular_sample(budget*0.9, budget, budget*1.1, N)

# 4) Automation learning productivity function (diminishing returns)
#    prod_gain% per employee = A * (1 - exp(-B * hours)) * mgmt * econ * noise
A = 18.0  # max plausible monthly productivity % from training focus on automation
B = 0.55  # curvature of diminishing returns
noise = rng.normal(1.0, 0.08, N)  # implementation variance
prod_gain_pct_emp = A * (1 - np.exp(-B * hours_per_emp)) * mgmt_mult * econ_mult * noise

# 5) Decision-making time reduction (%), influenced by mgmt & econ
base_decision_reduction = rng.uniform(decision_reduction_min, decision_reduction_max, N)
decision_reduction_pct = base_decision_reduction * mgmt_mult * econ_mult * rng.normal(1.0, 0.05, N)

# 6) Redundancy hours saved per employee per month (linked to prod gain)
redundant_hours_saved = baseline_redundant_hours * (prod_gain_pct_emp / 100.0) * rng.normal(1.0, 0.08, N)
redundant_hours_saved = np.clip(redundant_hours_saved, 0, None)

# 7) Satisfaction (eNPS-like) — centered around 75 with influence of support & hours
satisfaction = bounded_normal(mean=72 + 6*(mgmt_mult-1.0) + 2.0*(hours_per_emp-1.5),
                              sd=8, lower=55, upper=95, size=N)

# 8) Retention uplift: 2–5% of baseline turnover, scaled by support
retention_uplift_pct = rng.uniform(0.02, 0.05, N) * (0.9 + 0.4*(mgmt_mult-1.0))
# Effective turnover reduction:
turnover_reduction_pct_points = (baseline_turnover/100.0) * retention_uplift_pct

# ----------------------------
# Cost/Savings Calculations (monthly)
# ----------------------------
# Fully loaded cost
FHC = fully_loaded_hourly

# Redundancy savings
monthly_redundancy_savings = engaged_participants * redundant_hours_saved * FHC

# Decision-making savings
# total hours in decision processes per month
decision_hours_total = decision_volume * avg_hours_per_decision
monthly_decision_savings = decision_hours_total * (decision_reduction_pct/100.0) * FHC

# Retention savings (annualized, then convert to monthly)
replacement_cost = avg_annual_salary * replacement_cost_factor
employees_saved_per_year = engaged_participants * turnover_reduction_pct_points
annual_retention_savings = employees_saved_per_year * replacement_cost
monthly_retention_savings = annual_retention_savings / 12.0

# Total monthly savings
monthly_total_savings = monthly_redundancy_savings + monthly_decision_savings + monthly_retention_savings

# Productivity target check (per company level)
# Define company-level improvement as weighted average gain * engagement rate
company_prod_gain_pct = prod_gain_pct_emp * (engaged_participants / np.maximum(participants,1))
hit_target = company_prod_gain_pct >= prod_target

# ROI months to breakeven
# Guard against zero or tiny savings
roi_months = np.where(monthly_total_savings > 1, budget_spend / monthly_total_savings, np.inf)
roi_months = np.clip(roi_months, 0, 60)  # cap to 5 years for display

# ----------------------------
# KPIs (Executive)
# ----------------------------
prob_hit_target = np.mean(hit_target)
median_roi = float(np.nanmedian(roi_months))
p80_roi = float(np.nanpercentile(roi_months, 80))
p20_roi = float(np.nanpercentile(roi_months, 20))

median_monthly_savings = float(np.nanmedian(monthly_total_savings))
median_redundancy = float(np.nanmedian(monthly_redundancy_savings))
median_decision = float(np.nanmedian(monthly_decision_savings))
median_retention = float(np.nanmedian(monthly_retention_savings))
median_satisfaction = float(np.nanmedian(satisfaction))

# Decision Speed Index (extra metric): effective cycle time multiplier = 1 - decision_reduction%
decision_speed_index = 1 - np.nanmedian(decision_reduction_pct)/100.0

# Capability Maturity Lift (extra metric): normalized gain 0-1 vs max A
capability_lift = float(np.nanmedian(prod_gain_pct_emp))/A

# ----------------------------
# Top KPI Cards
# ----------------------------
def kpi(label, value, sublabel=None):
    st.markdown(f"""
    <div class="kpi-card">
      <div class="kpi-label">{label}</div>
      <div class="kpi-value">{value}</div>
      <div class="kpi-sub">{sublabel if sublabel else ""}</div>
    </div>
    """, unsafe_allow_html=True)

col1, col2, col3, col4 = st.columns(4)
with col1: kpi("Prob. Hit Productivity Target", f"{prob_hit_target*100:,.0f}%", f"Target ≥ {prod_target}%")
with col2: kpi("Median ROI (Months)", f"{median_roi:,.1f}", f"P20 {p20_roi:,.1f} • P80 {p80_roi:,.1f}")
with col3: kpi("Median Monthly Hard Savings", f"${median_monthly_savings:,.0f}",
               f"Redundancy ${median_redundancy:,.0f} • Decisions ${median_decision:,.0f}")
with col4: kpi("Employee Satisfaction (Median)", f"{median_satisfaction:,.0f}/100",
               f"Decision Speed Index {decision_speed_index:,.2f}")

st.markdown('<div class="small">Capability Maturity Lift: '
            f'<b>{capability_lift:.2f}</b> (0–1 scale)</div>', unsafe_allow_html=True)

# ----------------------------
# Distributions & Visuals
# ----------------------------
tab1, tab2, tab3 = st.tabs(["ROI timeline", "Savings breakdown", "Assumptions & details"])

with tab1:
    st.subheader("⏳ ROI Timeline Distribution (Months to Breakeven)")
    fig_roi = px.histogram(pd.DataFrame({"ROI (months)": roi_months}), x="ROI (months)", nbins=40,
                           color_discrete_sequence=["#0ea5e9"], opacity=0.85)
    fig_roi.add_vline(x=median_roi, line_width=2, line_dash="dash", line_color="#f59e0b")
    fig_roi.update_layout(height=420, bargap=0.05, margin=dict(l=0,r=0,t=30,b=0))
    st.plotly_chart(fig_roi, use_container_width=True)

    st.caption("Note: ROI months are capped at 60 for display. Infinite values (no breakeven) are extremely rare with current assumptions.")

with tab2:
    st.subheader("💸 Monthly Hard Savings Distribution")
    df_sav = pd.DataFrame({
        "Total": monthly_total_savings,
        "Redundancy": monthly_redundancy_savings,
        "Decision-making": monthly_decision_savings,
        "Retention (allocated monthly)": monthly_retention_savings
    })
    fig_total = px.histogram(df_sav, x="Total", nbins=40, color_discrete_sequence=["#22c55e"], opacity=0.85)
    fig_total.update_layout(height=320, margin=dict(l=0,r=0,t=30,b=0))
    st.plotly_chart(fig_total, use_container_width=True)

    c1, c2 = st.columns(2)
    with c1:
        fig_stack = go.Figure()
        for col in ["Redundancy", "Decision-making", "Retention (allocated monthly)"]:
            fig_stack.add_trace(go.Histogram(x=df_sav[col], nbinsx=40, name=col, opacity=0.6))
        fig_stack.update_layout(barmode='overlay', height=320, margin=dict(l=0,r=0,t=30,b=0))
        st.plotly_chart(fig_stack, use_container_width=True)
    with c2:
        st.markdown("**Median Monthly Savings Breakdown**")
        bd = pd.DataFrame({
            "Component": ["Redundancy", "Decision-making", "Retention (monthly)"],
            "Median $": [median_redundancy, median_decision, median_retention]
        })
        fig_bar = px.bar(bd, x="Component", y="Median $", text="Median $",
                         color="Component", color_discrete_sequence=["#10b981","#0ea5e9","#a78bfa"])
        fig_bar.update_traces(texttemplate='$%{text:,.0f}', textposition='outside')
        fig_bar.update_layout(yaxis_title=None, xaxis_title=None, height=320, margin=dict(l=0,r=0,t=30,b=0))
        st.plotly_chart(fig_bar, use_container_width=True)

with tab3:
    st.subheader("📋 Model Inputs & Assumptions")
    a1, a2 = st.columns(2)
    with a1:
        st.write("**Scenario & Context**")
        st.json({
            "Scenario": scenario if scenario != "Custom configuration" else "Custom",
            "Participants (planned)": participants,
            "Budget (planned)": budget,
            "Mgmt support": mgmt_support,
            "Economic condition": econ,
            "Hours per employee (planned)": hours
        })
        st.write("**Financial & Workforce**")
        st.json({
            "Fully loaded hourly cost": fully_loaded_hourly,
            "Baseline redundant hours/mo": baseline_redundant_hours,
            "Decision volume/mo": decision_volume,
            "Avg hours per decision": avg_hours_per_decision,
            "Baseline turnover % (annual)": baseline_turnover,
            "Average salary": avg_annual_salary,
            "Replacement cost factor": replacement_cost_factor
        })
    with a2:
        st.write("**Targets & Distributions**")
        st.json({
            "Productivity target %": prod_target,
            "Decision reduction range %": [decision_reduction_min, decision_reduction_max],
            "Model A (max prod gain %)": A,
            "Model B (curvature)": B,
            "Run count": N
        })
        st.write("**Derived medians**")
        st.json({
            "Engaged participants (median)": int(np.median(engaged_participants)),
            "Company productivity gain % (median)": float(np.median(company_prod_gain_pct)),
            "Decision reduction % (median)": float(np.median(decision_reduction_pct)),
            "Satisfaction (median)": median_satisfaction
        })

# ----------------------------
# Executive Summary Table (for PPT)
# ----------------------------
summary_rows = [
    ("Scenario", scenario if scenario != "Custom configuration" else "Custom"),
    ("Planned Participants", f"{participants:,}"),
    ("Planned Budget", f"${budget:,.0f}"),
    ("Mgmt Support", mgmt_support),
    ("Economic Condition", econ),
    ("Prob. Hit Productivity Target", f"{prob_hit_target*100:,.0f}% (Target ≥ {prod_target}%)"),
    ("Median ROI (Months)", f"{median_roi:,.1f} (P20 {p20_roi:,.1f}, P80 {p80_roi:,.1f})"),
    ("Median Monthly Total Savings", f"${median_monthly_savings:,.0f}"),
    ("Median Redundancy Savings", f"${median_redundancy:,.0f}"),
    ("Median Decision Savings", f"${median_decision:,.0f}"),
    ("Median Retention Savings (mo)", f"${median_retention:,.0f}"),
    ("Median Satisfaction", f"{median_satisfaction:,.0f}/100"),
    ("Decision Speed Index", f"{decision_speed_index:.2f} (lower is faster)"),
    ("Capability Maturity Lift", f"{capability_lift:.2f} (0–1)")
]
summary_df = pd.DataFrame(summary_rows, columns=["KPI", "Value"])

st.markdown("### 📑 Executive Summary")
st.dataframe(summary_df, use_container_width=True, hide_index=True)

# ----------------------------
# PPT Export
# ----------------------------
dl_col1, dl_col2 = st.columns([1,2])
with dl_col1:
    if st.button("📤 Export one-pager to PowerPoint"):
        pptx_bytes = export_to_ppt(summary_df, roi_months, scenario)
        st.download_button("Download PPTX", data=pptx_bytes, file_name="FutureReady_Simulation_Summary.pptx",
                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
with dl_col2:
    st.caption("Exports title, KPIs, and ROI distribution histogram — ready to drop into your board deck.")

# ----------------------------
# Bottom note
# ----------------------------
st.markdown("""
<div class="small">
Assumptions emphasize **hard savings** (redundancy, decision-making, retention) and a storyline view (Pilot, Conservative, Aggressive).
Tune sliders to match your context; distributions include realistic variance and management/economic multipliers.
</div>
""", unsafe_allow_html=True)
